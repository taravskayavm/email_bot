from __future__ import annotations

import io
import logging
import os
import re
import tempfile
import time
import zipfile
from pathlib import PurePosixPath
from typing import Any, Dict, Iterator, List, Optional, Tuple

from emailbot import settings
from emailbot.progress_watchdog import ProgressTracker, heartbeat_now
from emailbot.timebudget import TimeBudget
from .extraction_common import filter_invalid_tld
from .extraction_pdf import extract_text_from_pdf_bytes
from .reporting import log_extract_digest
from emailbot.utils import zip_limits as zl
from emailbot.utils.logging_setup import get_logger
from .utils.timeouts import DEFAULT_TIMEOUT_SEC, run_with_timeout

logger = get_logger(__name__)


_PROGRESS_KEYS = {"files_total", "files_processed", "files_skipped_timeout", "last_file"}


def _inspect_zip_members(path: str) -> List[Tuple[str, int]]:
    """Inspect ``path`` and return safe member names with their sizes."""

    max_files = int(getattr(zl, "MAX_FILES_PER_ZIP", 500))
    max_total = int(getattr(zl, "MAX_TOTAL_UNCOMPRESSED_BYTES", 200 * 1024 * 1024))
    with zipfile.ZipFile(path) as archive:
        members = archive.infolist()
        if len(members) > max_files:
            raise ValueError(f"ZIP limit: files={len(members)} > {max_files}")
        total_size = 0
        safe_members: List[Tuple[str, int]] = []
        for info in members:
            if info.is_dir():
                continue
            total_size += int(getattr(info, "file_size", 0) or 0)
            if total_size > max_total:
                raise ValueError("ZIP limit: total uncompressed size exceeded")
            normalized = os.path.normpath(info.filename)
            if normalized.startswith(".."):
                continue
            if normalized.startswith(("/", "\\")):
                continue
            safe_members.append((normalized, int(getattr(info, "file_size", 0) or 0)))
        return safe_members


def extract_from_zip(path: str, timeout_sec: int = DEFAULT_TIMEOUT_SEC) -> List[str]:
    """Return a sanitized list of member names for ``path``.

    The inspection is executed inside :func:`run_with_timeout` to prevent hangs on
    crafted archives.  When a timeout or validation failure occurs an empty list
    is returned.
    """

    try:
        members = run_with_timeout(_inspect_zip_members, timeout_sec, path)
    except TimeoutError:
        logger.warning("ZIP inspect timeout: %s after %ss", path, timeout_sec)
        return []
    except (ValueError, zipfile.BadZipFile) as exc:
        logger.warning("ZIP blocked %s: %s", path, exc)
        return []
    return [name for name, _size in members]


ZIP_MAX_FILES = int(os.getenv("ZIP_MAX_FILES", "500"))
ZIP_MAX_TOTAL_UNCOMP_MB = int(os.getenv("ZIP_MAX_TOTAL_UNCOMP_MB", "500"))
ZIP_MAX_MEMBER_MB = int(os.getenv("ZIP_MAX_MEMBER_MB", "50"))
ZIP_MAX_DEPTH = int(os.getenv("ZIP_MAX_DEPTH", "2"))
ZIP_RATIO_LIMIT = float(os.getenv("ZIP_RATIO_LIMIT", "100.0"))
# Сколько секунд даём на обработку одного элемента архива (PDF/DOCX/и т.п.)
# Увеличено по умолчанию до 60 из-за «ложных зависаний» на тяжёлых документах.
ZIP_MEMBER_TIMEOUT_SEC = int(os.getenv("ZIP_MEMBER_TIMEOUT_SEC", "60"))
_XLSX_CELL_LIMIT = int(os.getenv("ZIP_XLSX_CELL_LIMIT", "5000"))
_PPTX_TEXT_LIMIT = int(os.getenv("ZIP_PPTX_TEXT_LIMIT", "2000"))


def _is_suspicious(member: zipfile.ZipInfo) -> bool:
    comp = max(member.compress_size or 1, 1)
    uncomp = max(member.file_size or 1, 1)
    ratio = float(uncomp) / float(comp)
    return (ratio > ZIP_RATIO_LIMIT) or (
        uncomp > ZIP_MAX_MEMBER_MB * 1024 * 1024
    )


def _decode_text(data: bytes) -> str:
    for encoding in ("utf-8", "utf-16", "cp1251", "latin-1"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", "ignore")


def _odt_to_text(data: bytes) -> str:
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as odt:
            xml_bytes = odt.read("content.xml")
    except Exception:
        return ""
    text = _decode_text(xml_bytes)
    text = re.sub(r"<[^>]+>", " ", text)
    return re.sub(r"\s+", " ", text).strip()


def extract_text_from_bytes_guess(
    name: str, data: bytes, budget: TimeBudget | None = None
) -> str:
    """Attempt to extract plain text from ``data`` based on ``name`` extension."""

    ext = os.path.splitext(name or "")[1].lower()
    heartbeat_now()
    if budget:
        budget.checkpoint()

    if ext == ".pdf":
        return extract_text_from_pdf_bytes(data, budget=budget)
    if ext in {".txt", ".csv", ".html", ".htm", ".md", ".json"}:
        return _decode_text(data)
    if ext in {".docx", ".doc"}:
        try:
            import docx  # type: ignore
        except Exception:
            return ""
        try:
            doc = docx.Document(io.BytesIO(data))
        except Exception:
            return ""
        parts: list[str] = []
        for para in doc.paragraphs:
            heartbeat_now()
            if budget:
                budget.checkpoint()
            if para.text:
                parts.append(para.text)
        return "\n".join(parts)
    if ext == ".rtf":
        try:
            from striprtf.striprtf import rtf_to_text  # type: ignore
        except Exception:
            return _decode_text(data)
        try:
            return rtf_to_text(_decode_text(data))
        except Exception:
            return _decode_text(data)
    if ext == ".odt":
        return _odt_to_text(data)
    if ext == ".xlsx":
        try:
            import openpyxl  # type: ignore
        except Exception:
            return ""
        try:
            wb = openpyxl.load_workbook(
                io.BytesIO(data), read_only=True, data_only=True
            )
        except Exception:
            return ""
        texts: list[str] = []
        cells = 0
        try:
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    heartbeat_now()
                    if budget:
                        budget.checkpoint()
                    for value in row:
                        if value is None:
                            continue
                        texts.append(str(value))
                        cells += 1
                        if cells >= _XLSX_CELL_LIMIT:
                            break
                    if cells >= _XLSX_CELL_LIMIT:
                        break
                if cells >= _XLSX_CELL_LIMIT:
                    break
        finally:
            try:
                wb.close()
            except Exception:
                pass
        return "\n".join(texts)
    if ext == ".pptx":
        try:
            from pptx import Presentation  # type: ignore
        except Exception:
            return ""
        try:
            presentation = Presentation(io.BytesIO(data))
        except Exception:
            return ""
        texts: list[str] = []
        for slide in presentation.slides:
            heartbeat_now()
            if budget:
                budget.checkpoint()
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                texts.append(shape.text)
                if len(texts) >= _PPTX_TEXT_LIMIT:
                    break
            if len(texts) >= _PPTX_TEXT_LIMIT:
                break
        return "\n".join(texts)
    return ""


def _extract_text_from_zipfile(
    zf: zipfile.ZipFile, budget: TimeBudget | None, depth: int
) -> str:
    out_chunks: list[str] = []
    total_uncomp = 0
    files_seen = 0
    for info in zf.infolist():
        heartbeat_now()
        if budget:
            budget.checkpoint()
        if files_seen >= ZIP_MAX_FILES:
            break
        if info.is_dir():
            continue
        if info.flag_bits & 0x1:
            continue
        name = info.filename or ""
        if not _safe_path(name):
            continue
        if _is_suspicious(info):
            continue
        files_seen += 1
        read_limit = min(info.file_size or 0, ZIP_MAX_MEMBER_MB * 1024 * 1024)
        try:
            with zf.open(info, "r") as source:
                payload = source.read(read_limit)
        except Exception:
            continue
        total_uncomp += len(payload)
        if total_uncomp > ZIP_MAX_TOTAL_UNCOMP_MB * 1024 * 1024:
            break
        lower_name = name.lower()
        if lower_name.endswith(".zip") and depth + 1 <= ZIP_MAX_DEPTH:
            try:
                with zipfile.ZipFile(io.BytesIO(payload), "r") as nested:
                    nested_text = _extract_text_from_zipfile(
                        nested, budget, depth + 1
                    )
            except Exception:
                nested_text = ""
            if nested_text:
                out_chunks.append("\n" + nested_text)
            continue
        text = extract_text_from_bytes_guess(lower_name, payload, budget=budget)
        if text:
            out_chunks.append("\n" + text)
    return "".join(out_chunks)


def extract_text_from_zip(
    path: str, budget: TimeBudget | None = None, depth: int = 0
) -> str:
    if depth > ZIP_MAX_DEPTH:
        return ""
    try:
        with zipfile.ZipFile(path, "r") as zf:
            return _extract_text_from_zipfile(zf, budget, depth)
    except Exception:
        return ""


def _safe_unlink(path: str, attempts: int = 6, delay: float = 0.2) -> bool:
    import os, time, logging

    log = logging.getLogger(__name__)
    for i in range(attempts):
        try:
            os.remove(path)
            return True
        except PermissionError:
            time.sleep(delay * (i + 1))
        except FileNotFoundError:
            return True
        except Exception:
            log.exception("unlink failed for %s", path)
            break
    return False

ALLOWED_EXTS = {".pdf", ".docx", ".xlsx", ".csv", ".txt", ".html", ".htm", ".zip"}
DENY_EXTS = {
    ".exe",
    ".dll",
    ".js",
    ".bat",
    ".cmd",
    ".sh",
    ".com",
    ".pif",
    ".scr",
    ".cpl",
    ".msi",
    ".msp",
    ".jar",
    ".vbs",
    ".ps1",
    ".php",
}

MAX_FILES = 1000
MAX_SIZE = 100 * 1024 * 1024  # 100 MB
MAX_DEPTH = 3


class _ZipLimitError(RuntimeError):
    """Raised when a ZIP archive violates protective limits."""

    def __init__(self, reason: str, *, total_size: int | None = None) -> None:
        super().__init__(reason)
        self.reason = reason
        self.total_size = total_size


def _iter_zip_member_names(z: zipfile.ZipFile) -> Iterator[str]:
    """Yield file member names while enforcing archive limits."""

    total_size = 0
    for idx, info in enumerate(z.infolist(), 1):
        if idx > MAX_FILES:
            raise _ZipLimitError("too many files")
        file_size = getattr(info, "file_size", 0) or 0
        total_size += file_size
        if total_size > MAX_SIZE:
            raise _ZipLimitError("too big", total_size=total_size)
        if info.is_dir():
            continue
        yield info.filename or ""


def _safe_path(name: str) -> bool:
    p = PurePosixPath(name.replace("\\", "/"))
    return not (p.is_absolute() or ".." in p.parts)


def extract_emails_from_zip(
    path: str,
    stop_event: Optional[object] = None,
    *,
    _depth: int = 0,
    tracker: ProgressTracker | None = None,
) -> tuple[list["EmailHit"], Dict[str, Any]]:
    from .extraction import (
        EmailHit,
        extract_any_stream,
        merge_footnote_prefix_variants,
        repair_footnote_singletons,
        _dedupe,
    )

    start = time.monotonic()
    if _depth > MAX_DEPTH:
        logger.warning("zip depth exceeded for %s", path)
        return [], {"errors": ["max depth exceeded"]}

    try:
        z = zipfile.ZipFile(path)
    except Exception:
        return [], {"errors": ["cannot open"]}

    hits: List[EmailHit] = []
    stats: Dict[str, Any] = {
        "files_total": 0,
        "files_processed": 0,
        "files_skipped_timeout": 0,
        "last_file": "",
    }

    def _tracker_update(**kwargs: Any) -> None:
        if tracker is None:
            return
        updater = getattr(tracker, "update", None)
        if not callable(updater):
            return
        payload = dict(kwargs)
        done_value = payload.get("done")
        processed_value = payload.get("processed")
        if done_value is not None and processed_value is None:
            payload["processed"] = done_value
        elif processed_value is not None and done_value is None:
            payload["done"] = processed_value
        elif processed_value is None:
            current_processed = int(stats.get("files_processed", 0))
            payload["processed"] = current_processed
            payload.setdefault("done", current_processed)
        if "file" not in payload and payload.get("current") is not None:
            payload["file"] = payload["current"]
        if "total" not in payload:
            payload["total"] = int(stats.get("files_total", 0))
        try:
            updater(**payload)
        except Exception:
            logger.debug("tracker update failed", exc_info=True)

    try:
        member_names = list(_iter_zip_member_names(z))
        total_members = len(member_names)
        stats["files_total"] = total_members
        if tracker is not None and _depth == 0:
            tracker.reset_total(total_members)
            _tracker_update(stage="scan", done=0, total=total_members)
        elif _depth == 0:
            _tracker_update(stage="scan", done=0, total=total_members)
        if _depth == 0 and total_members:
            _tracker_update(
                stage="process",
                current=None,
                done=int(stats.get("files_processed", 0)),
                total=total_members,
            )
        for name in member_names:
            if stop_event and getattr(stop_event, "is_set", lambda: False)():
                break
            stats["last_file"] = name
            heartbeat_now()
            file_start = time.monotonic()
            last_pulse = file_start
            processed_tick = False
            if _depth == 0:
                _tracker_update(
                    stage="prepare",
                    current=name,
                    done=int(stats.get("files_processed", 0)),
                    total=total_members,
                )

            def _pulse(stage: str, *, force: bool = False) -> None:
                nonlocal last_pulse
                now = time.monotonic()
                if force or (now - last_pulse) >= 0.3:
                    heartbeat_now()
                    if _depth == 0:
                        _tracker_update(
                            stage=stage,
                            current=name,
                            done=int(stats.get("files_processed", 0)),
                            total=total_members,
                            elapsed=round(now - file_start, 3),
                        )
                    last_pulse = now

            try:
                try:
                    info = z.getinfo(name)
                except KeyError:
                    logger.warning("missing zip entry in %s: %s", path, name)
                    stats["zip_member_error"] = stats.get("zip_member_error", 0) + 1
                    continue

                _pulse("processing", force=True)

                if info.flag_bits & 0x1:
                    logger.warning("encrypted file skipped in zip %s: %s", path, name)
                    continue
                if not _safe_path(name):
                    logger.warning("unsafe path in zip %s: %s", path, name)
                    return [], {"errors": ["unsafe path"]}
                ext = os.path.splitext(name)[1].lower()
                if ext in DENY_EXTS:
                    logger.warning("deny-listed extension in zip %s: %s", path, name)
                    return [], {"errors": ["forbidden extension"]}
                if ext == ".zip":
                    if _depth + 1 > MAX_DEPTH:
                        logger.warning("nested zip depth exceeded in %s: %s", path, name)
                        return [], {"errors": ["max depth exceeded"]}
                    data = z.read(info)
                    _pulse("processing")
                    with tempfile.NamedTemporaryFile(delete=False, suffix=".zip") as tmp:
                        tmp.write(data)
                        tmp_path = tmp.name
                    _pulse("processing")
                    try:
                        inner_hits, inner_stats = extract_emails_from_zip(
                            tmp_path,
                            stop_event,
                            _depth=_depth + 1,
                            tracker=tracker,
                        )
                        for h in inner_hits:
                            suffix = ""
                            if "#" in h.source_ref:
                                suffix = "#" + h.source_ref.split("#", 1)[1]
                            new_ref = f"zip:{path}|{name}{suffix}"
                            hits.append(
                                EmailHit(
                                    email=h.email,
                                    source_ref=new_ref,
                                    origin=h.origin,
                                    pre=h.pre,
                                    post=h.post,
                                )
                            )
                        for k, v in inner_stats.items():
                            if k in _PROGRESS_KEYS:
                                continue
                            if isinstance(v, int):
                                stats[k] = stats.get(k, 0) + v
                        current_total = int(stats.get("files_total", 0))
                        inner_total = int(inner_stats.get("files_total") or 0)
                        stats["files_total"] = max(current_total + inner_total - 1, 0)
                        if tracker is not None and inner_total > 0:
                            tracker.extend_total(max(inner_total - 1, 0))
                        stats["files_processed"] = int(stats.get("files_processed", 0)) + int(
                            inner_stats.get("files_processed") or 0
                        )
                        stats["files_skipped_timeout"] = int(
                            stats.get("files_skipped_timeout", 0)
                        ) + int(inner_stats.get("files_skipped_timeout") or 0)
                        last_inner = inner_stats.get("last_file")
                        if last_inner:
                            stats["last_file"] = last_inner
                        processed_tick = True
                    finally:
                        if not _safe_unlink(tmp_path):
                            logger.warning("temp file still locked, skip delete: %s", tmp_path)
                    continue
                if ext not in ALLOWED_EXTS:
                    continue
                data = z.read(info)
                _pulse("processing")
                source_ref = f"zip:{path}|{name}"

                try:
                    inner_hits, inner_stats = run_with_timeout(
                        extract_any_stream,
                        ZIP_MEMBER_TIMEOUT_SEC,
                        data,
                        ext,
                        source_ref=source_ref,
                        stop_event=stop_event,
                    )
                    _pulse("processing")
                except TimeoutError:
                    logger.warning(
                        "zip member timed out",
                        extra={"event": "zip_member_timeout", "entry": source_ref},
                    )
                    stats["zip_member_timeout"] = stats.get("zip_member_timeout", 0) + 1
                    stats["files_skipped_timeout"] = int(
                        stats.get("files_skipped_timeout", 0)
                    ) + 1
                    continue
                except Exception as exc:  # pragma: no cover - defensive logging
                    logger.debug(
                        "zip member extraction error",
                        extra={
                            "event": "zip_member_error",
                            "entry": source_ref,
                            "error": repr(exc),
                        },
                    )
                    stats["zip_member_error"] = stats.get("zip_member_error", 0) + 1
                    continue

                hits.extend(inner_hits)
                key = ext.lstrip(".")
                stats[key] = stats.get(key, 0) + 1
                for k, v in inner_stats.items():
                    if k in _PROGRESS_KEYS:
                        continue
                    if isinstance(v, int):
                        stats[k] = stats.get(k, 0) + v
                stats["files_processed"] = int(stats.get("files_processed", 0)) + 1
                processed_tick = True
                if _depth == 0:
                    _pulse("processing", force=True)
            finally:
                if tracker is not None:
                    tracker.tick_file(name, processed=processed_tick)
                if _depth == 0:
                    elapsed = round(time.monotonic() - file_start, 3)
                    _tracker_update(
                        stage="done_file",
                        current=name,
                        done=int(stats.get("files_processed", 0)),
                        total=total_members,
                        elapsed=elapsed,
                    )
        if _depth == 0:
            _tracker_update(
                stage="done",
                done=int(stats.get("files_processed", 0)),
                total=int(stats.get("files_total", 0)),
            )
    except _ZipLimitError as exc:
        if exc.reason == "too many files":
            logger.warning("zip has too many files: %s", path)
        elif exc.reason == "too big":
            total_size = exc.total_size or 0
            logger.warning("zip too large: %s (%d bytes)", path, total_size)
        return [], {"errors": [exc.reason]}
    finally:
        try:
            z.close()
        except Exception:
            pass
    hits = merge_footnote_prefix_variants(hits, stats)
    hits, fstats = repair_footnote_singletons(hits, settings.PDF_LAYOUT_AWARE)
    for k, v in fstats.items():
        if v:
            stats[k] = stats.get(k, 0) + v
    hits = _dedupe(hits)
    emails, extra = filter_invalid_tld([h.email for h in hits], stats=stats)
    stats["invalid_tld"] = stats.get("invalid_tld", 0) + extra.get("invalid_tld", 0)
    replacements = extra.get("replacements") or {}
    if replacements:
        updated_hits: list[EmailHit] = []
        for h in hits:
            new_email = replacements.get(h.email)
            if new_email:
                updated_hits.append(
                    EmailHit(
                        email=new_email,
                        source_ref=h.source_ref,
                        origin=h.origin,
                        pre=h.pre,
                        post=h.post,
                        meta=h.meta,
                    )
                )
            else:
                updated_hits.append(h)
        hits = _dedupe(updated_hits)
    samples = extra.get("invalid_tld_examples") or []
    if samples:
        stored = stats.setdefault("invalid_tld_examples", [])
        for sample in samples:
            if sample not in stored:
                stored.append(sample)
            if len(stored) >= 3:
                break
    allowed = set(emails)
    hits = [h for h in hits if h.email in allowed]
    stats["unique_after_cleanup"] = len(hits)
    suspicious = sum(1 for h in hits if h.email.split("@", 1)[0].isdigit())
    if suspicious:
        stats["suspicious_numeric_localpart"] = stats.get(
            "suspicious_numeric_localpart", 0
        ) + suspicious
    stats["mode"] = "file"
    stats["entry"] = path
    stats["elapsed_ms"] = int((time.monotonic() - start) * 1000)
    log_extract_digest(stats)
    return hits, stats


__all__ = [
    "extract_from_zip",
    "extract_emails_from_zip",
    "extract_text_from_zip",
    "extract_text_from_bytes_guess",
]
